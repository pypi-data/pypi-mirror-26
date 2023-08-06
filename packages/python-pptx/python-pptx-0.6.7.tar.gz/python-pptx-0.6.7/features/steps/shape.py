# encoding: utf-8

"""
Gherkin step implementations for shape-related features.
"""

from __future__ import absolute_import, print_function

import hashlib

from behave import given, when, then

from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import (
    MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE, PP_MEDIA_TYPE
)
from pptx.action import ActionSetting
from pptx.util import Emu, Inches

from helpers import (
    cls_qname, saved_pptx_path, test_file, test_pptx, test_text
)


# given ===================================================

@given('an autoshape')
def given_an_autoshape(context):
    prs = Presentation(test_pptx('shp-autoshape-adjustments'))
    context.shape = prs.slides[0].shapes[0]


@given('an autoshape having text')
def given_an_autoshape_having_text(context):
    prs = Presentation(test_pptx('shp-autoshape-props'))
    context.shape = prs.slides[0].shapes[0]


@given("(builder._start_x, builder._start_y) is ({x_str}, {y_str})")
def given_builder_start_x_builder_start_y_is_x_y(context, x_str, y_str):
    builder = context.builder
    builder._start_x, builder._start_y = int(x_str), int(y_str)


@given("(builder._x_scale, builder._y_scale) is ({p_str}, {q_str})")
def given_builder_x_scale_builder_y_scale_is_p_q(context, p_str, q_str):
    builder = context.builder
    builder._x_scale, builder._y_scale = float(p_str), float(q_str)


@given('a chevron shape')
def given_a_chevron_shape(context):
    prs = Presentation(test_pptx('shp-autoshape-adjustments'))
    context.shape = prs.slides[0].shapes[0]


@given('a connector')
def given_a_connector(context):
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[4]


@given('a connector and a 1 inch square picture at 0, 0')
def given_a_connector_and_a_1_inch_square_picture_at_0_0(context):
    prs = Presentation(test_pptx('shp-connector-props'))
    shapes = prs.slides[1].shapes
    context.picture = shapes[0]
    context.connector = shapes[1]


@given('a connector having its begin point at ({x}, {y})')
def given_a_connector_having_its_begin_point_at_x_y(context, x, y):
    prs = Presentation(test_pptx('shp-connector-props'))
    sld = prs.slides[0]
    context.connector = sld.shapes[0]


@given('a connector having its end point at ({x}, {y})')
def given_a_connector_having_its_end_point_at_x_y(context, x, y):
    prs = Presentation(test_pptx('shp-connector-props'))
    sld = prs.slides[0]
    context.connector = sld.shapes[0]


@given('a FreeformBuilder object as builder')
def given_a_FreeformBuilder_object_as_builder(context):
    shapes = Presentation(test_pptx('shp-freeform')).slides[0].shapes
    builder = shapes.build_freeform()
    context.builder = builder


@given('a graphic frame')  # shouldn't matter, but this one contains a table
def given_a_graphic_frame(context):
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[2]


@given('a graphic frame containing a chart')
def given_a_graphic_frame_containing_a_chart(context):
    prs = Presentation(test_pptx('shp-access-chart'))
    sld = prs.slides[0]
    context.shape = sld.shapes[0]


@given('a graphic frame containing a table')
def given_a_graphic_frame_containing_a_table(context):
    prs = Presentation(test_pptx('shp-access-chart'))
    sld = prs.slides[1]
    context.shape = sld.shapes[0]


@given('a group shape')
def given_a_group_shape(context):
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[3]


@given('a movie shape')
def given_a_movie_shape(context):
    prs = Presentation(test_pptx('shp-movie-props'))
    context.movie = prs.slides[0].shapes[0]


@given('a picture')
def given_a_picture(context):
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[1]


@given('a rotated {shape_type}')
def given_a_rotated_shape_type(context, shape_type):
    shape_idx = {
        'shape':         0,
        'picture':       1,
        'graphic frame': 2,
        'group shape':   3,
        'connector':     4,
    }[shape_type]
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[1]
    context.shape = sld.shapes[shape_idx]


@given('a shape')
def given_a_shape(context):
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[0]


@given('a SlideShapes object as shapes')
def given_a_SlideShapes_object_as_shapes(context):
    prs = Presentation(test_pptx('shp-shape-access'))
    context.shapes = prs.slides[0].shapes


@given('a SlideShapes object containing {a_or_no} movies')
def given_a_SlideShapes_object_containing_a_or_no_movies(context, a_or_no):
    pptx = {
        'one or more': 'shp-movie-props',
        'no':          'shp-shape-access',
    }[a_or_no]
    prs = Presentation(test_pptx(pptx))
    context.prs = prs
    context.shapes = prs.slides[0].shapes


@given('a SlideShapes object containing 6 shapes')
def given_a_SlideShapes_object_containing_6_shapes(context):
    prs = Presentation(test_pptx('shp-shape-access'))
    context.shapes = prs.slides[0].shapes


@given('a SlideShapes object having a {type} shape at offset {idx}')
def given_a_SlideShapes_obj_having_type_shape_at_off_idx(context, type, idx):
    prs = Presentation(test_pptx('shp-shape-access'))
    context.shapes = prs.slides[1].shapes


@given('a textbox')
def given_a_textbox(context):
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[5]


@given('a {shape_type} on a slide')
def given_a_shape_on_a_slide(context, shape_type):
    shape_idx = {
        'shape':         0,
        'picture':       1,
        'graphic frame': 2,
        'group shape':   3,
        'connector':     4,
    }[shape_type]
    prs = Presentation(test_pptx('shp-common-props'))
    sld = prs.slides[0]
    context.shape = sld.shapes[shape_idx]
    context.slide = sld


@given('a shape of known position and size')
def given_a_shape_of_known_position_and_size(context):
    prs = Presentation(test_pptx('shp-pos-and-size'))
    context.shape = prs.slides[0].shapes[0]


# when ====================================================

@when("I add a text box to the slide's shape collection")
def when_I_add_a_text_box(context):
    shapes = context.slide.shapes
    x, y = (Inches(1.00), Inches(2.00))
    cx, cy = (Inches(3.00), Inches(1.00))
    sp = shapes.add_textbox(x, y, cx, cy)
    sp.text = test_text


@when("I add an auto shape to the slide's shape collection")
def when_I_add_an_auto_shape(context):
    shapes = context.slide.shapes
    x, y = (Inches(1.00), Inches(2.00))
    cx, cy = (Inches(3.00), Inches(4.00))
    sp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    sp.text = test_text


@when("I assign 0.15 to shape.adjustments[0]")
def when_I_assign_to_shape_adjustments(context):
    context.shape.adjustments[0] = 0.15


@when("I assign a string to shape.text")
def when_I_assign_a_string_to_shape_text(context):
    context.shape.text = u'F\xf8o\nBar'


@when("I assign builder.convert_to_shape() to shape")
def when_I_assign_builder_convert_to_shape_to_shape(context):
    builder = context.builder
    context.shape = builder.convert_to_shape()


@when("I assign builder.convert_to_shape({x_str}, {y_str}) to shape")
def when_I_assign_builder_convert_to_shape_origin_x_y(context, x_str, y_str):
    builder = context.builder
    origin_x, origin_y = int(x_str), int(y_str)
    context.shape = builder.convert_to_shape(origin_x, origin_y)


@when("I assign shapes.build_freeform() to builder")
def when_I_assign_shapes_build_freeform_to_builder(context):
    shapes = context.shapes
    builder = shapes.build_freeform()
    context.builder = builder


@when("I assign shapes.build_freeform(scale=100.0) to builder")
def when_I_assign_shapes_build_freeform_scale_to_builder(context):
    shapes = context.shapes
    builder = shapes.build_freeform(scale=100.0)
    context.builder = builder


@when("I assign shapes.build_freeform(scale=(200.0, 100.0)) to builder")
def when_I_assign_shapes_build_freeform_scale_rectnglr_to_builder(context):
    shapes = context.shapes
    builder = shapes.build_freeform(scale=(200.0, 100.0))
    context.builder = builder


@when("I assign shapes.build_freeform(start_x=25, start_y=125) to builder")
def when_I_assign_shapes_build_freeform_start_x_start_y_to_builder(context):
    shapes = context.shapes
    builder = shapes.build_freeform(25, 125)
    context.builder = builder


@when('I assign {value} to connector.begin_x')
def when_I_assign_value_to_connector_begin_x(context, value):
    context.connector.begin_x = int(value)


@when('I assign {value} to connector.begin_y')
def when_I_assign_value_to_connector_begin_y(context, value):
    context.connector.begin_y = int(value)


@when('I assign {value} to connector.end_x')
def when_I_assign_value_to_connector_end_x(context, value):
    context.connector.end_x = int(value)


@when('I assign {value} to connector.end_y')
def when_I_assign_value_to_connector_end_y(context, value):
    context.connector.end_y = int(value)


@when("I assign '{value}' to shape.name")
def when_I_assign_value_to_shape_name(context, value):
    context.shape.name = value


@when("I assign {value} to shape.rotation")
def when_I_assign_value_to_shape_rotation(context, value):
    context.shape.rotation = float(value)


@when('I call builder.add_line_segments([(100, 25), (25, 100)])')
def when_I_call_builder_add_line_segments_100_25_25_100(context):
    builder = context.builder
    builder.add_line_segments([(100, 25), (25, 100)])


@when('I call connector.begin_connect(picture, 3)')
def when_I_call_connector_begin_connect_picture_3(context):
    connector, picture = context.connector, context.picture
    connector.begin_connect(picture, 3)


@when('I call connector.end_connect(picture, 3)')
def when_I_call_connector_end_connect_picture_3(context):
    connector, picture = context.connector, context.picture
    connector.end_connect(picture, 3)


@when('I call shapes.add_chart({type_}, chart_data)')
def when_I_call_shapes_add_chart(context, type_):
    chart_type = getattr(XL_CHART_TYPE, type_)
    context.chart = context.shapes.add_chart(
        chart_type, 0, 0, 0, 0, context.chart_data
    ).chart


@when('I call shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 1, 2, 3, 4)')
def when_I_call_shapes_add_connector(context):
    context.connector = context.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 1, 2, 3, 4
    )


@when('I call shapes.add_movie(file, x, y, cx, cy, poster_frame)')
def when_I_call_shapes_add_movie(context):
    shapes = context.shapes
    x, y, cx, cy = Emu(2590800), Emu(571500), Emu(3962400), Emu(5715000)
    context.movie = shapes.add_movie(
        test_file('just-two-mice.mp4'), x, y, cx, cy,
        test_file('just-two-mice.png')
    )


@when("I change the left and top of the {shape_type}")
def when_I_change_the_position_of_the_shape(context, shape_type):
    left, top = {
        'shape':          (692696, 1339552),
        'picture':       (1835696, 2711152),
        'graphic frame': (2978696, 4082752),
        'group shape':   (4121696, 5454352),
        'connector':     (5264696, 6825952),
    }[shape_type]
    shape = context.shape
    shape.left = left
    shape.top = top


@when("I change the width and height of the {shape_type}")
def when_I_change_the_size_of_the_shape(context, shape_type):
    width, height = {
        'shape':          (692696, 1339552),
        'picture':       (1835696, 2711152),
        'graphic frame': (2978696, 4082752),
        'group shape':   (4121696, 5454352),
        'connector':     (5264696, 6825952),
    }[shape_type]
    shape = context.shape
    shape.width = width
    shape.height = height


@when("I get the chart from its graphic frame")
def when_I_get_the_chart_from_its_graphic_frame(context):
    context.chart = context.shape.chart


# then ====================================================

@then("builder is a FreeformBuilder object")
def then_builder_is_a_FreeformBuilder_object(context):
    builder = context.builder
    class_name = builder.__class__.__name__
    expected_value = 'FreeformBuilder'
    assert class_name == expected_value, (
        'Expected class name \'%s\', got \'%s\'' %
        (expected_value, class_name)
    )


@then("(builder._start_x, builder._start_y) is ({x_str}, {y_str})")
def then_builder_start_x_builder_start_y_is_x_y(context, x_str, y_str):
    builder = context.builder
    actual_value = builder._start_x, builder._start_y
    expected_value = int(x_str), int(y_str)
    assert actual_value == expected_value, (
        'Expected %s, got %s' % (expected_value, actual_value)
    )


@then("(builder._x_scale, builder._y_scale) is ({p_str}, {q_str})")
def then_builder_x_scale_builder_y_scale_is_x_y(context, p_str, q_str):
    builder = context.builder
    actual_value = builder._x_scale, builder._y_scale
    expected_value = float(p_str), float(q_str)
    assert actual_value == expected_value, (
        'Expected %s, got %s' % (expected_value, actual_value)
    )


@then('connector is a Connector object')
def then_connector_is_a_Connector_object(context):
    assert type(context.connector).__name__ == 'Connector'


@then('connector.begin_x == {value}')
def then_connector_begin_x_equals_value(context, value):
    assert context.connector.begin_x == int(value)


@then('connector.begin_x is an Emu object with value {x}')
def then_connector_begin_x_is_an_Emu_object_with_value_x(context, x):
    begin_x = context.connector.begin_x
    assert isinstance(begin_x, Emu)
    assert begin_x == int(x)


@then('connector.begin_y == {value}')
def then_connector_begin_y_equals_value(context, value):
    assert context.connector.begin_y == int(value)


@then('connector.begin_y is an Emu object with value {y}')
def then_connector_begin_y_is_an_Emu_object_with_value_y(context, y):
    begin_y = context.connector.begin_y
    assert isinstance(begin_y, Emu)
    assert begin_y == int(y)


@then('connector.end_x == {value}')
def then_connector_end_x_equals_value(context, value):
    assert context.connector.end_x == int(value)


@then('connector.end_x is an Emu object with value {x}')
def then_connector_end_x_is_an_Emu_object_with_value_x(context, x):
    end_x = context.connector.end_x
    assert isinstance(end_x, Emu)
    assert end_x == int(x)


@then('connector.end_y == {value}')
def then_connector_end_y_equals_value(context, value):
    assert context.connector.end_y == int(value)


@then('connector.end_y is an Emu object with value {y}')
def then_connector_end_y_is_an_Emu_object_with_value_y(context, y):
    end_y = context.connector.end_y
    assert isinstance(end_y, Emu)
    assert end_y == int(y)


@then('movie is a Movie object')
def then_movie_is_a_Movie_object(context):
    class_name = context.movie.__class__.__name__
    assert class_name == 'Movie', 'got %s' % class_name


@then("movie.left, movie.top == x, y")
def then_movie_left_movie_top_eq_x_y(context):
    movie = context.movie
    position = movie.left, movie.top
    assert position == (Emu(2590800), Emu(571500)), 'got %s' % position


@then('movie.media_format is a _MediaFormat object')
def then_movie_media_format_is_a_MediaFormat_object(context):
    class_name = context.movie.media_format.__class__.__name__
    assert class_name == '_MediaFormat', 'got %s' % class_name


@then('movie.media_type is PP_MEDIA_TYPE.MOVIE')
def then_movie_media_type_is_PP_MEDIA_TYPE_MOVIE(context):
    media_type = context.movie.media_type
    assert media_type == PP_MEDIA_TYPE.MOVIE, 'got %s' % media_type


@then("movie.poster_frame is the same image as poster_frame")
def then_movie_poster_frame_is_the_same_image_as_poster_frame(context):
    actual_sha1 = context.movie.poster_frame.sha1
    with open(test_file('just-two-mice.png'), 'rb') as f:
        expected_sha1 = hashlib.sha1(f.read()).hexdigest()
    assert actual_sha1 == expected_sha1, 'not the same image'


@then('movie.shape_type is MSO_SHAPE_TYPE.MEDIA')
def then_movie_shape_type_is_MSO_SHAPE_TYPE_MEDIA(context):
    shape_type = context.movie.shape_type
    assert shape_type == MSO_SHAPE_TYPE.MEDIA, 'got %s' % shape_type


@then("movie.width, movie.height == cx, cy")
def then_movie_width_movie_height_eq_cx_cy(context):
    movie = context.movie
    size = movie.width, movie.height
    assert size == (Emu(3962400), Emu(5715000)), 'got %s' % size


@then('shape.adjustments[0] is 0.15')
def then_shape_adjustments_is_value(context):
    shape = context.shape
    assert shape.adjustments[0] == 0.15


@then("shape.click_action is an ActionSetting object")
def then_shape_click_action_is_an_ActionSetting_object(context):
    assert isinstance(context.shape.click_action, ActionSetting)


@then('shape.has_text_frame is {value_str}')
def then_shape_has_text_frame_is(context, value_str):
    expected_value = {'True': True, 'False': False}[value_str]
    has_text_frame = context.shape.has_text_frame
    assert has_text_frame is expected_value, 'got %s' % has_text_frame


@then("(shape.left, shape.top) is ({x_str}, {y_str})")
def then_shape_left_shape_top_is_x_y(context, x_str, y_str):
    shape = context.shape
    actual_value = shape.left, shape.top
    expected_value = int(x_str), int(y_str)
    assert actual_value == expected_value, (
        'Expected %s, got %s' % (expected_value, actual_value)
    )


@then('shape.line is a LineFormat object')
def then_shape_line_is_a_LineFormat_object(context):
    shape = context.shape
    line_format = shape.line
    line_format_cls_name = cls_qname(line_format)
    expected_cls_name = 'pptx.dml.line.LineFormat'
    assert line_format_cls_name == expected_cls_name, (
        "expected '%s', got '%s'" % (expected_cls_name, line_format_cls_name)
    )


@then("shape.name is '{expected_value}'")
def then_shape_name_is_value(context, expected_value):
    shape = context.shape
    msg = "expected shape name '%s', got '%s'" % (shape.name, expected_value)
    assert shape.name == expected_value, msg


@then('shape.part is the SlidePart of the shape')
def then_shape_part_is_the_SlidePart_of_the_shape(context):
    assert context.shape.part is context.slide.part


@then("shape.rotation is {value}")
def then_shape_rotation_is_value(context, value):
    shape = context.shape
    expected_value = float(value)
    assert shape.rotation == expected_value, 'got %s' % expected_value


@then('shape.shape_id == {value_str}')
def then_shape_shape_id_equals(context, value_str):
    expected_value = int(value_str)
    shape_id = context.shape.shape_id
    assert shape_id == expected_value, 'got %s' % shape_id


@then('shape.text is the string I assigned')
def then_shape_text_is_the_string_I_assigned(context):
    shape = context.shape
    assert shape.text == u'F\xf8o\nBar'


@then('shape.text is the text in the shape')
def then_shape_text_is_the_text_in_the_shape(context):
    shape = context.shape
    assert shape.text == u'Fee Fi\nF\xf8\xf8 Fum\nI am a shape\nwith textium'


@then("(shape.width, shape.height) is ({cx_str}, {cy_str})")
def then_shape_width_shape_height_is_cx_cy(context, cx_str, cy_str):
    shape = context.shape
    actual_value = shape.width, shape.height
    expected_value = int(cx_str), int(cy_str)
    assert actual_value == expected_value, (
        'Expected %s, got %s' % (expected_value, actual_value)
    )


@then('the auto shape appears in the slide')
def then_auto_shape_appears_in_slide(context):
    prs = Presentation(saved_pptx_path)
    sp = prs.slides[0].shapes[0]
    sp_text = sp.text_frame.paragraphs[0].runs[0].text
    assert sp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    assert sp.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    assert sp_text == test_text


@then('the chart is a Chart object')
def then_the_chart_is_a_Chart_object(context):
    assert isinstance(context.chart, Chart)


@then('the left and top of the {shape_type} match their new values')
def then_left_and_top_of_shape_match_new_values(context, shape_type):
    expected_left, expected_top = {
        'shape':          (692696, 1339552),
        'picture':       (1835696, 2711152),
        'graphic frame': (2978696, 4082752),
        'group shape':   (4121696, 5454352),
        'connector':     (5264696, 6825952),
    }[shape_type]
    shape = context.shape
    assert shape.left == expected_left, 'got left: %s' % shape.left
    assert shape.top == expected_top, 'got top: %s' % shape.top


@then('the left and top of the {shape_type} match their known values')
def then_left_and_top_of_shape_match_known_values(context, shape_type):
    expected_left, expected_top = {
        'shape':         (1339552,  692696),
        'picture':       (2711152, 1835696),
        'graphic frame': (4082752, 2978696),
        'group shape':   (5454352, 4121696),
        'connector':     (6825952, 5264696),
    }[shape_type]
    shape = context.shape
    assert shape.left == expected_left, 'got left: %s' % shape.left
    assert shape.top == expected_top, 'got top: %s' % shape.top


@then('the shape {has_or_not} a chart')
def then_the_shape_has_or_not_a_chart(context, has_or_not):
    expected_bool = {'has': True, 'does not have': False}[has_or_not]
    shape = context.shape
    assert shape.has_chart is expected_bool


@then('the width and height of the {shape_type} match their known values')
def then_width_and_height_of_shape_match_known_values(context, shape_type):
    expected_width, expected_height = {
        'shape':         (928192, 914400),
        'picture':       (914400, 945232),
        'graphic frame': (993304, 914400),
        'group shape':   (914400, 914400),
        'connector':     (986408, 828600),
    }[shape_type]
    shape = context.shape
    assert shape.width == expected_width, 'got width: %s' % shape.width
    assert shape.height == expected_height, 'got height: %s' % shape.height


@then('the width and height of the {shape_type} match their new values')
def then_width_and_height_of_shape_match_new_values(context, shape_type):
    expected_width, expected_height = {
        'shape':          (692696, 1339552),
        'picture':       (1835696, 2711152),
        'graphic frame': (2978696, 4082752),
        'group shape':   (4121696, 5454352),
        'connector':     (5264696, 6825952),
    }[shape_type]
    shape = context.shape
    assert shape.width == expected_width, 'got width: %s' % shape.width
    assert shape.height == expected_height, 'got height: %s' % shape.height


@then('the text box appears in the slide')
def then_text_box_appears_in_slide(context):
    prs = Presentation(saved_pptx_path)
    textbox = prs.slides[0].shapes[0]
    textbox_text = textbox.text_frame.paragraphs[0].runs[0].text
    assert textbox_text == test_text
