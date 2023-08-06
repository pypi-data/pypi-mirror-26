"""Python-PPTX customized module."""
import os
import six
import copy
import logging
import requests
import datetime
import tempfile
import collections
import numpy as np
from . import utils
import pandas as pd
import matplotlib.cm
from . import fontwidth
import matplotlib.colors
from . import color as _color
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from collections import OrderedDict
from orderedattrdict import AttrDict
from tornado.template import Template
from tornado.escape import to_unicode
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import XyChartData
from pptx.oxml.xmlchemy import OxmlElement
from pptx.chart.data import BubbleChartData
from six.moves.urllib_parse import urlparse
from gramex.transforms import build_transform

_template_cache = {}


def template(tmpl, data):
    """Execute tornado template."""
    if tmpl not in _template_cache:
        _template_cache[tmpl] = Template(tmpl, autoescape=None)
    return to_unicode(_template_cache[tmpl].generate(**data))


def text(shape, spec, data):
    '''
    Replace entire text of shape with spec['text'].
    '''
    if not shape.has_text_frame:
        logging.error('"%s" is not a TextShape to apply text:', shape.name)
        return
    if not isinstance(data, (dict,)):
        data = {'data': data}
    run_flag = True
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = template(spec['text'], data) if run_flag else ''
            run_flag = False


def replace(shape, spec, data):
    '''
    Replace keywords in shape using the dictionary at spec['replace'].
    '''
    if not shape.has_text_frame:
        logging.error('"%s" is not a TextShape to apply text:', shape.name)
        return
    if not isinstance(data, (dict,)):
        data = {'data': data}
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for old, new in spec['replace'].items():
                run.text = run.text.replace(old, template(new, data))


def image(shape, spec, data):
    '''
    Replace image with a different file specified in spec['image']
    '''
    image = template(spec['image'], data)
    # If it's a URL, use the requests library's raw stream as a file-like object
    if urlparse(image).netloc:
        r = requests.get(image)
        with tempfile.NamedTemporaryFile(delete=False) as handle:
            handle.write(r.content)
        new_img_part, new_rid = shape.part.get_or_add_image_part(handle.name)
        os.unlink(handle.name)
    else:
        new_img_part, new_rid = shape.part.get_or_add_image_part(image)
    old_rid = shape._pic.blip_rId
    shape._pic.blipFill.blip.rEmbed = new_rid
    shape.part.related_parts[old_rid].blob = new_img_part.blob


def stack_elements(replica, shape, stack=False, margin=None):
    '''
    Function to extend elements horizontally or vertically.
    '''
    config = {'vertical': {'axis': 'y', 'attr': 'height'},
              'horizontal': {'axis': 'x', 'attr': 'width'}}
    if stack is not None:
        grp_sp = shape.element
        # Adding a 15% margin between original and new object.
        margin = 0.50 if not margin else margin
        for index in range(replica - 1):
            # Adding a cloned object to shape
            extend_shape = copy.deepcopy(grp_sp)
            # Getting attributes and axis values from config based on stack.
            attr = config.get(stack, {}).get('attr', 0)
            axis = config.get(stack, {}).get('axis', 0)
            # Taking width or height based on stack value and setting a margin.
            metric_val = getattr(shape, attr)
            axis_val = getattr(extend_shape, axis)
            # Setting margin accordingly either vertically or horizontally.
            axis_pos = metric_val * (index + 1)
            set_attr = axis_val + axis_pos + (axis_pos * margin)
            # Setting graphic position of newly created object to slide.
            setattr(extend_shape, axis, int(set_attr))
            # Adding newly created object to slide.
            grp_sp.addnext(extend_shape)


def stack_shapes(collection, change, data):
    '''
    Function to stack Shapes if required.
    '''
    shape_data = data
    for shape in collection:
        if shape.name not in change:
            continue
        info = change[shape.name]
        if 'data' in info and info.get('stack') is not None:
            shape_data = eval(info.get('data'))
        stack_elements(len(shape_data), shape, stack=info.get('stack'),
                       margin=info.get('margin'))


def generate_slide(prs, source):
    '''
    Create a slide layout.
    '''
    layout_items_count = [
        len(layout.placeholders) for layout in prs.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return prs.slide_layouts[blank_layout_id]


def copy_slide_elem(shape, dest):
    '''
    Function to copy slide elements into a newly created slide.
    '''
    if dest:
        el = shape.element
        new_elem = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(new_elem, 'p:extLst')


def add_new_slide(dest, source_slide):
    '''
    Function to add a new slide to presentation.
    '''
    if dest:
        for key, value in six.iteritems(source_slide.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                dest.part.rels.add_relationship(value.reltype, value._target, value.rId)


def replicate_slide(change, prs, data):
    '''
    Add Duplicate slide to the presentation.
    '''
    if 'replicate-slides' in change and change['replicate-slides'] is not None:
        replica_conf = change['replicate-slides']['slide-number']
        for conf in replica_conf:
            for slide_number, replica_logic in conf.items():
                replica = template(replica_logic, data)
                for replica in range(int(replica) - 1):
                    source = prs.slides[slide_number - 1]
                    add_new_slide(prs, source)


def move_slide(presentation, old_index, new_index):
    '''
    Move a slide's index number.
    '''
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def delete_slide(presentation, index):
    '''
    Delete a slide from Presentation.
    '''
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


def rect_css(shape, **kwargs):
    """Function to add text to shape."""
    for key in {'fill', 'stroke'}:
        if kwargs.get(key):
            fill = shape.fill if key == 'fill' else shape.line.fill
            rectcss = kwargs[key].rsplit('#')[-1].lower()
            rectcss = rectcss + ('0' * (6 - len(rectcss)))
            chart_css(fill, kwargs, rectcss)


def add_text_to_shape(shape, textval, **kwargs):
    """Function to add text to shape."""
    min_inc = 13000
    pixel_inch = 10000
    # kwargs['font-size'] = max(kwargs.get('font-size', 16), min_inc)
    if (kwargs.get('font-size', 14) * pixel_inch) > min_inc:
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.add_run()
        for run in paragraph.runs:
            run.text = textval
            shape_txt = run.font
            shape_txt = run.font.fill
            shape_txt.solid()
            utils.apply_text_css(run, paragraph, **kwargs)


def scale_data(data, lo, hi, factor=None):
    """Function to scale data."""
    return ((data - lo) / (hi - lo)) * factor if factor else (data - lo) / (hi - lo)


def rect(shape, x, y, width, height):
    """Add rectangle to slide."""
    return shape.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)


def _update_chart(info, data, chart_data, series_columns, chart='line'):
    """Updating Chart data."""
    series_dict = {}
    columns = list(data.columns)
    if info['x'] in columns:
        columns.remove(info['x'])

    if chart == 'line':
        for series in series_columns:
            if np.issubdtype(data[series].dtype, np.number):
                chart_data.add_series(series, tuple(data[series].fillna(0).tolist()))
        return chart_data

    for index, row in data.fillna(0).iterrows():
        for col in series_columns:
            if col in columns and np.issubdtype(data[col].dtype, np.number):
                serieslist = [series.name for series in chart_data._series]
                if col not in serieslist:
                    series_dict[col] = chart_data.add_series(col)
                try:
                    x = float(row[info['x']])
                except Exception:
                    x = index + 1
                if chart == 'scatter':
                    series_dict[col].add_data_point(x, row[col])
                elif chart == 'bubble':
                    bubble_size = row[info['size']] if info.get('size') else 1
                    if col != info.get('size'):
                        series_dict[col].add_data_point(x, row[col], bubble_size)
    return chart_data


def chart_css(fill, info, color):
    """Function to add opacity to charts."""
    fill.solid()
    pix_to_inch = 100000
    fill.fore_color.rgb = RGBColor.from_string(color)
    solid_fill = fill.fore_color._xFill
    alpha = OxmlElement('a:alpha')
    alpha.set('val', '%d' % (pix_to_inch * info.get('opacity', 1.0)))
    solid_fill.srgbClr.append(alpha)
    return fill


def compile_function(spec, key, data, handler):
    """A function to compile configuration."""
    if not spec.get(key):
        return None
    _vars = {'_color': None, 'data': None, 'handler': None}
    if not isinstance(spec[key], (dict,)):
        spec[key] = {'function': '{}'.format(spec[key])}
    elif isinstance(spec[key], (dict,)) and 'function' not in spec[key]:
        spec[key] = {'function': '{}'.format(spec[key])}
    args = {'data': data, 'handler': handler, '_color': _color}
    return build_transform(spec[key], vars=_vars)(**args)[0]


def table(shape, spec, data):
    """Update an existing Table shape with data."""
    if not spec.get('table', {}).get('data'):
        return
    spec = copy.deepcopy(spec['table'])
    handler = data.pop('handler') if 'handler' in data else None
    data = compile_function(spec, 'data', data, handler)
    if not len(data):
        return
    data_cols = list(data.columns)
    data_cols = list(set(spec.get('columns', {}).keys() or data_cols).intersection(data_cols))
    table_properties = utils.TableProperties()
    # Extending table if required.
    table_properties.extend_table(shape, data, len(data) + 1, len(data_cols))
    # Fetching Table Style for All Cells and texts.
    tbl_style = table_properties.get_default_css(shape)
    cell_style = table_properties.get_css(spec, data_cols, data)
    data = data.to_dict(orient='records')
    for row_num, row in enumerate(shape.table.rows):
        cols = len(row.cells._tr.tc_lst)
        # Extending cells in newly added rows.
        while cols < len(data_cols):
            row.cells._tr.add_tc()
            cols += 1
        for col_num, cell in enumerate(row.cells):
            colname = data_cols[col_num]
            for paragraph in cell.text_frame.paragraphs:
                if not paragraph.text.strip():
                    paragraph.add_run()
                for run in paragraph.runs:
                    txt = colname if row_num == 0 else data[row_num - 1][colname]
                    run.text = '{}'.format(txt)
                    cellcss = {} if row_num == 0 else copy.deepcopy(cell_style.get(colname, {}))
                    txt_css = copy.deepcopy(tbl_style.get('header' if row_num == 0 else 'row', {}))
                    if row_num > 0 and 'gradient' in cellcss:
                        grad_txt = scale_data(txt, cellcss['min'], cellcss['max'])
                        gradient = matplotlib.cm.get_cmap(cellcss['gradient'])
                        cellcss['fill'] = matplotlib.colors.to_hex(gradient(grad_txt))
                    if cellcss.get('fill'):
                        txt_css['color'] = cellcss.get('color', _color.contrast(cellcss['fill']))
                        cellcss['color'] = cellcss.get('color', _color.contrast(cellcss['fill']))
                    txt_css.update(cellcss)
                    table_properties.apply_table_css(cell, paragraph, run, txt_css)


def chart(shape, spec, data):
    """Replacing chart Data."""
    chart_type = None
    if hasattr(shape.chart, 'chart_type'):
        chart_type = '{}'.format(shape.chart.chart_type).split()[0]

    stacked_or_line = {
        'AREA', 'AREA_STACKED', 'AREA_STACKED_100', 'BAR_CLUSTERED',
        'BAR_OF_PIE', 'BAR_STACKED', 'BAR_STACKED_100', 'COLUMN_CLUSTERED',
        'COLUMN_STACKED', 'COLUMN_STACKED_100', 'LINE',
        'LINE_MARKERS', 'LINE_MARKERS_STACKED', 'LINE_MARKERS_STACKED_100',
        'LINE_STACKED', 'LINE_STACKED_100', 'RADAR_MARKERS',
        'RADAR', 'RADAR_FILLED', 'PIE', 'PIE_EXPLODED', 'PIE_OF_PIE',
        'DOUGHNUT', 'DOUGHNUT_EXPLODED'}

    xy_charts = {
        'XY_SCATTER', 'XY_SCATTER_LINES', 'XY_SCATTER_LINES_NO_MARKERS',
        'XY_SCATTER_SMOOTH', 'XY_SCATTER_SMOOTH_NO_MARKERS'}

    bubble_charts = {'BUBBLE', 'BUBBLE_THREE_D_EFFECT'}
    if not chart_type:
        raise NotImplementedError()

    info = copy.deepcopy(spec['chart'])
    # Load data
    handler = data.pop('handler') if 'handler' in data else None
    if isinstance(info['x'], (dict,)):
        if 'function' not in info['x']:
            info['x']['function'] = '{}'.format(info['x'])
        info['x'] = compile_function(info, 'x', data, handler)

    if 'size' in info and isinstance(info['size'], (dict,)):
        if 'function' not in info['size']:
            info['size']['function'] = '{}'.format(info['size'])
        info['size'] = compile_function(info, 'size', data, handler)
    data = compile_function(info, 'data', data, handler)
    series_cols = [x for x in data.columns if x != info['x']]
    chart_name = None
    # If chart type is stacked bar or line.
    if chart_type in stacked_or_line:
        # Initializing chart data
        chart_name = 'line'
        chart_data = ChartData()
        chart_data.categories = data[info['x']].dropna().unique().tolist()
        data = _update_chart(info, data, chart_data, series_cols)

    # If chart type is scatter plot.
    elif chart_type in xy_charts:
        # Initializing chart data
        chart_name = 'scatter'
        chart_data = XyChartData()
        data = _update_chart(info, data, chart_data, series_cols, chart='scatter')

    # If chart type is bubble chart.
    elif chart_type in bubble_charts:
        # Initializing chart data
        chart_name = 'bubble'
        chart_data = BubbleChartData()
        data = _update_chart(info, data, chart_data, series_cols, chart='bubble')
    else:
        raise NotImplementedError()

    shape.chart.replace_data(chart_data)

    if 'color' in info and isinstance(info['color'], (dict,)):
        if 'function' not in info['color']:
            info['color']['function'] = '{}'.format(info['color'])
        info['color'] = compile_function(info, 'color', data, handler)

    if 'opacity' in info and isinstance(info['opacity'], (dict,)):
        if 'function' not in info['opacity']:
            info['opacity']['function'] = '{}'.format(info['opacity'])
        info['opacity'] = compile_function(info, 'opacity', data, handler)

    if chart_name == 'scatter' and not info.get('color'):
        series_names = [series.name for series in shape.chart.series]
        info['color'] = dict(zip(series_names, _color.distinct(len(series_names))))

    if info.get('color'):
        for index, series in enumerate(shape.chart.series):
            fill_graph = utils.conver_color_code(info['color'].get(series.name, '#cccccc'))
            if chart_name == 'scatter':
                fill = series.marker.format.fill
                line_fill = series.marker.format.line.fill
            else:
                fill = series.format.fill
                line_fill = series.format.line.fill
            chart_css(fill, info, fill_graph)
            chart_css(line_fill, info, fill_graph)

# Custom Charts Functions below(Sankey, Treemap, Calendarmap).


def sankey(shape, spec, data):
    """Draw sankey in PPT."""
    # Shape must be a rectangle.
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()
    # Getting parent shapes
    pxl_to_inch = 10000
    default_thickness = 40
    spec = copy.deepcopy(spec['sankey'])

    handler = data.pop('handler') if 'handler' in data else None
    shapes = shape._parent
    y0 = shape.top
    x0 = shape.left

    width = shape.width
    shape_ids = {'shape': 0}
    height = shape.height

    groups = compile_function(spec, 'groups', data, handler)
    thickness = spec.get('thickness', default_thickness) * pxl_to_inch

    h = (height - (thickness * len(groups))) / (len(groups) - 1) + thickness
    frames = {}
    # Sankey Rectangles and texts.
    sankey_conf = {}
    sankey_conf['x0'] = x0
    sankey_conf['size'] = compile_function(spec, 'size', data, handler)
    sankey_conf['width'] = width
    sankey_conf['order'] = compile_function(spec, 'order', data, handler)
    sankey_conf['text'] = compile_function(spec, 'text', data, handler)
    sankey_conf['color'] = compile_function(spec, 'color', data, handler)
    sankey_conf['attrs'] = spec.get('attrs', {})
    sankey_conf['sort'] = spec.get('sort')
    stroke = spec.get('stroke', '#ffffff')
    # Delete rectangle after geting width, height, x-position and y-position
    shape._sp.delete()
    elem_schema = utils.make_element()
    data = compile_function(spec, 'data', data, handler)
    for ibar, group in enumerate(groups):
        y = y0 + h * ibar
        sankey_conf['group'] = [group]
        df = frames[group] = utils.draw_sankey(data, sankey_conf)
        # Adding rectangle
        for key, row in df.iterrows():
            shp = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, row['x'], y, row['width'], thickness)
            rectstyle = {"fill": row['fill'], 'stroke': stroke}
            rect_css(shp, **rectstyle)
            text_style = {}
            text_style['color'] = _color.contrast(row['fill'])
            text_style.update(spec.get('style', {}))
            add_text_to_shape(shp, row['text'], **text_style)

    # Sankey Connection Arcs.
    for ibar, (group1, group2) in enumerate(zip(groups[:-1], groups[1:])):
        sankey_conf['group'] = [group1, group2]
        sankey_conf['sort'] = False
        df = utils.draw_sankey(data, sankey_conf)
        pos = collections.defaultdict(float)
        for key1, row1 in frames[group1].iterrows():
            for key2, row2 in frames[group2].iterrows():
                if (key1, key2) in df.index:
                    row = df.ix[(key1, key2)]
                    y1, y2 = y0 + h * ibar + thickness, y0 + h * (ibar + 1)
                    ym = (y1 + y2) / 2
                    x1 = row1['x'] + pos[0, key1]
                    x2 = row2['x'] + pos[1, key2]

                    _id = shape_ids['shape'] = shape_ids['shape'] + 1
                    shp = utils.cust_shape(
                        0, 0, '{:.0f}'.format(row['width']), '{:.0f}'.format(ym), _id)
                    path = elem_schema['a'].path(
                        w='{:.0f}'.format(row['width']), h='{:.0f}'.format(ym))
                    shp.find('.//a:custGeom', namespaces=elem_schema['nsmap']).append(
                        elem_schema['a'].pathLst(path))
                    path.append(
                        elem_schema['a'].moveTo(elem_schema['a'].pt(
                            x='{:.0f}'.format(x1 + row['width']), y='{:.0f}'.format(y1))))

                    path.append(elem_schema['a'].cubicBezTo(
                        elem_schema['a'].pt(x='{:.0f}'.format(x1 + row['width']),
                                            y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x2 + row['width']),
                                            y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x2 + row['width']),
                                            y='{:.0f}'.format(y2))))

                    path.append(elem_schema['a'].lnTo(
                        elem_schema['a'].pt(x='{:.0f}'.format(x2), y='{:.0f}'.format(y2))))

                    path.append(elem_schema['a'].cubicBezTo(
                        elem_schema['a'].pt(x='{:.0f}'.format(x2), y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x1), y='{:.0f}'.format(ym)),
                        elem_schema['a'].pt(x='{:.0f}'.format(x1), y='{:.0f}'.format(y1))))

                    path.append(elem_schema['a'].close())
                    shp.spPr.append(elem_schema['a'].solidFill(
                        utils.fill_color(srgbclr=row['fill'])))
                    shapes._spTree.append(shp)
                    pos[0, key1] += row['width']
                    pos[1, key2] += row['width']


def treemap(shape, spec, data):
    """Function to download data as ppt."""
    # Shape must be a rectangle.
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()
    shapes = shape._parent
    x0 = shape.left
    y0 = shape.top
    width = shape.width
    height = shape.height
    spec = copy.deepcopy(spec['treemap'])
    stroke = spec.get('stroke', '#ffffff')
    # Load data
    handler = data.pop('handler') if 'handler' in data else None
    spec['keys'] = compile_function(spec, 'keys', data, handler)
    spec['values'] = compile_function(spec, 'values', data, handler)
    spec['size'] = compile_function(spec, 'size', data, handler)
    spec['sort'] = compile_function(spec, 'sort', data, handler)
    spec['color'] = compile_function(spec, 'color', data, handler)
    spec['text'] = compile_function(spec, 'text', data, handler)
    spec['data'] = compile_function(spec, 'data', data, handler)
    # Getting rectangle's width and height using `squarified` algorithm.
    treemap_data = utils.SubTreemap(**spec)
    # Delete rectangle after geting width, height, x-position and y-position
    shape._sp.delete()
    font_aspect = 14.5
    pixel_inch = 10000
    default_rect_color = '#cccccc'
    for x, y, w, h, (l, v) in treemap_data.draw(width, height):
        if l == 0:
            shp = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x + x0, y + y0, w, h)
            if spec.get('color'):
                rect_color = spec['color'](v) if callable(spec['color']) else spec['color']
            else:
                default_rect_color
            if spec.get('text'):
                text = spec['text'](v) if callable(spec['text']) else spec['text']
            else:
                text = '{}'.format(v[1])
            rectstyle = {"fill": rect_color, 'stroke': stroke}
            rect_css(shp, **rectstyle)
            font_size = min(h, w * font_aspect / fontwidth.fontwidth('{}'.format(text)), pd.np.Inf)
            text_style = {}
            text_style['color'] = _color.contrast(rect_color)
            text_style.update(spec.get('style', {}))
            text_style['font-size'] = font_size / pixel_inch
            # Adding text inside rectangles
            add_text_to_shape(shp, text, **text_style)


def calendarmap(shape, spec, data):
    """Draw calendar map in PPT."""
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
        raise NotImplementedError()

    shapes = shape._parent
    spec = copy.deepcopy(spec['calendarmap'])
    handler = data.get('handler')
    # Load data
    data = compile_function(spec, 'data', data, handler)
    startdate = compile_function(spec, 'startdate', data, handler)

    pixel_inch = 10000
    size = spec.get('size', None)

    label_top = spec.get('label_top', 0) * pixel_inch
    label_left = spec.get('label_left', 0) * pixel_inch

    width = spec['width'] * pixel_inch
    shape_top = label_top + shape.top
    shape_left = label_left + shape.left
    y0 = width + shape_top
    x0 = width + shape_left

    # Deleting the shape
    shape.element.delete()
    # Style
    default_color = '#ffffff'
    default_line_color = '#787C74'
    default_txt_color = '#000000'

    style = spec.get('style', {})
    for key, value in style.items():
        if isinstance(value, (dict,)) and 'function' in value:
            style[key] = compile_function(style, key, data, handler)

    font_size = style.get('font-size', 12)
    stroke = style.get('stroke', '#ffffff')
    fill_rect = style.get('fill', '#cccccc')
    text_color = style.get('color', '#000000')
    # Treat infinities as nans when calculating scale
    scaledata = pd.Series(data).replace([pd.np.inf, -pd.np.inf], pd.np.nan)
    for key in {'lo', 'hi', 'weekstart'}:
        if isinstance(spec.get(key), (dict,)) and 'function' in spec.get(key):
            spec[key] = compile_function(spec, key, data, handler)

    lo_data = spec.get('lo', scaledata.min())
    range_data = spec.get('hi', scaledata.max()) - lo_data
    gradient = matplotlib.cm.get_cmap(spec.get('gradient', 'RdYlGn'))
    color = style.get('fill', lambda v: matplotlib.colors.to_hex(
        gradient((float(v) - lo_data) / range_data)) if not pd.isnull(v) else default_color)

    startweekday = (startdate.weekday() - spec.get('weekstart', 0)) % 7
    # Weekday Mean and format
    weekday_mean = pd.Series(
        [scaledata[(x - startweekday) % 7::7].mean() for x in range(7)])
    weekday_format = spec.get('format', '{:,.%df}' % utils.decimals(weekday_mean.values))
    # Weekly Mean and format
    weekly_mean = pd.Series([scaledata[max(0, x):x + 7].mean()
                             for x in range(-startweekday, len(scaledata) + 7, 7)])
    weekly_format = spec.get('format', '{:,.%df}' % utils.decimals(weekly_mean.values))

    # Scale sizes as square roots from 0 to max (not lowest to max -- these
    # should be an absolute scale)
    sizes = width * utils.scale(
        [v ** .5 for v in size], lo=0) if size is not None else [width] * len(scaledata)

    for i, val in enumerate(data):
        nx = (i + startweekday) // 7
        ny = (i + startweekday) % 7
        d = startdate + datetime.timedelta(days=i)
        fill = '#ffffff'
        if not pd.isnull(val):
            fill = color(val) if callable(color) else color

        shp = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x0 + (width * nx) + (width - sizes[i]) / 2,
            y0 + (width * ny) + (width - sizes[i]) / 2,
            sizes[i], sizes[i])
        rectstyle = {"fill": fill, 'stroke': stroke(val) if callable(stroke) else stroke}
        rect_css(shp, **rectstyle)
        text_style = {}
        text_style['color'] = style.get('color')(val) if callable(
            style.get('color')) else spec.get('color', _color.contrast(fill))
        text_style['font-size'] = font_size(val) if callable(font_size) else font_size
        text_style['bold'] = style.get('bold')
        text_style['italic'] = style.get('italic')
        text_style['underline'] = style.get('underline')
        text_style['font-family'] = style.get('font-family')
        add_text_to_shape(shp, '%02d' % d.day, **text_style)

        # Draw the boundary lines between months
        if i >= 7 and d.day == 1 and ny > 0:
            vertical_line = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x0 + width * nx, y0 + (width * ny), width, 2 * pixel_inch)
            rectstyle = {"fill": default_line_color, 'stroke': default_line_color}
            rect_css(vertical_line, **rectstyle)

        if i >= 7 and d.day <= 7 and nx > 0:
            horizontal_line = shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x0 + (width * nx), y0 + (width * ny), 2 * pixel_inch, width)
            rectstyle = {"fill": default_line_color, 'stroke': default_line_color}
            rect_css(horizontal_line, **rectstyle)

        # Adding weekdays text to the chart (left side)
        if i < 7:
            txt = shapes.add_textbox(
                x0 - (width / 2), y0 + (width * ny) + (width / 2), width, width)
            text_style['color'] = default_txt_color
            add_text_to_shape(txt, d.strftime('%a')[0], **text_style)

        # Adding months text to the chart (top)
        if d.day <= 7 and ny == 0:
            txt = shapes.add_textbox(
                x0 + (width * nx), y0 - (width / 2), width, width)
            text_style['color'] = default_txt_color
            add_text_to_shape(txt, d.strftime('%b %Y'), **text_style)

    if label_top:
        lo_weekly = spec.get('lo', weekly_mean.min())
        range_weekly = spec.get('hi', weekly_mean.max()) - lo_weekly
        for nx, val in enumerate(weekly_mean):
            if not pd.isnull(val):
                w = label_top * ((val - lo_weekly) / range_weekly)
                px = x0 + (width * nx)
                top_bar = shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, px, shape_top - w, width, w)
                rectstyle = {"fill": fill_rect(val) if callable(fill_rect) else fill_rect,
                             'stroke': stroke(val) if callable(stroke) else stroke}
                rect_css(top_bar, **rectstyle)
                top_txt = shapes.add_textbox(px, shape_top - width, width, width)
                text_style['color'] = text_color(val) if callable(text_color) else text_color
                add_text_to_shape(top_txt, weekly_format.format(weekly_mean[nx]), **text_style)

    if label_left:
        lo_weekday = spec.get('lo', weekday_mean.min())
        range_weekday = spec.get('hi', weekday_mean.max()) - lo_weekday
        for ny, val in enumerate(weekday_mean):
            if not pd.isnull(val):
                w = label_left * ((val - lo_weekday) / range_weekday)
                bar = shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, shape_left - w, y0 + (width * ny), w, width)
                rectstyle = {"fill": fill_rect(val) if callable(fill_rect) else fill_rect,
                             'stroke': stroke(val) if callable(stroke) else stroke}
                rect_css(bar, **rectstyle)
                left_txt = shapes.add_textbox(shape_left - width, y0 + (width * ny), w, width)
                text_style['color'] = text_color(val) if callable(text_color) else text_color
                add_text_to_shape(left_txt, weekday_format.format(weekday_mean[ny]), **text_style)


def bullet(shape, spec, data):
    """Function to plot bullet chart."""
    if shape.auto_shape_type != MSO_SHAPE.RECTANGLE:
            raise NotImplementedError()
    spec = copy.deepcopy(spec['bullet'])
    orient = spec.get('orient', 'horizontal')
    if orient not in {'horizontal', 'vertical'}:
            raise NotImplementedError()

    x = shape.left
    y = shape.top
    text = spec.get('text')
    height = shape.height if orient == 'horizontal' else shape.width
    width = shape.width if orient == 'horizontal' else shape.height
    handler = data.get('handler')
    for met in ['poor', 'average', 'good', 'target']:
        spec[met] = compile_function(spec, met, data, handler) if spec.get(met) else np.nan
    if text:
        text = compile_function({'text': spec['text']}, 'text', data, handler)

    spec['data'] = compile_function(spec, 'data', data, handler)
    gradient = spec.get('gradient', 'RdYlGn')
    shapes = shape._parent
    shape._sp.delete()
    lo = spec.get('lo', 0)
    hi = spec.get('hi', np.nanmax([spec['data'], spec['target'], spec['poor'],
                                  spec['average'], spec['good']]))
    default_font = 18
    style = {}
    common_style = copy.deepcopy(spec.get('style', {}))
    css = {'data': common_style.pop('data', {}), 'target': common_style.pop('target', {}),
           'poor': common_style.pop('poor', {}), 'good': common_style.pop('good', {}),
           'average': common_style.pop('average', {})}

    for key, val in css.items():
        _style = copy.deepcopy(common_style)
        _style.update(val)
        if 'font-size' not in _style:
            _style['font-size'] = default_font
        for css_prop, css_val in _style.items():
            if isinstance(css_val, (dict,)) and 'function' in css_val:
                _style[css_prop] = compile_function(_style, css_prop, data, handler)
        style[key] = _style

    gradient = matplotlib.cm.get_cmap(gradient)
    percentage = {'good': 0.125, 'average': 0.25, 'poor': 0.50}
    for index, metric in enumerate(['good', 'average', 'poor']):
        if not np.isnan(spec[metric]):
            scaled = scale_data(spec.get(metric, np.nan), lo, hi, factor=width)
            _width = scaled if orient == 'horizontal' else height
            _hight = height if orient == 'horizontal' else scaled
            yaxis = y if orient == 'horizontal' else y + (width - scaled)
            _rect = rect(shapes, x, yaxis, _width, _hight)
            fill = style.get(metric, {})
            stroke = fill.get('stroke')
            fill = fill.get('fill', matplotlib.colors.to_hex(gradient(percentage[metric])))
            rect_css(_rect, **{'fill': fill, 'stroke': stroke or fill})

    scaled = scale_data(spec['data'], lo, hi, factor=width)
    _width = scaled if orient == 'horizontal' else height / 2.0
    yaxis = y + height / 4.0 if orient == 'horizontal' else y + (width - scaled)
    xaxis = x if orient == 'horizontal' else x + height / 4.0
    _hight = height / 2.0 if orient == 'horizontal' else scaled
    data_rect = rect(shapes, xaxis, yaxis, _width, _hight)
    fill = style.get('data', {})
    stroke = fill.get('stroke')
    fill = fill.get('fill', matplotlib.colors.to_hex(gradient(1.0)))
    rect_css(data_rect, **{'fill': fill, 'stroke': stroke or fill})

    if text:
        data_text = text(spec['data'])
        parent = data_rect._parent
        _factor = 1.25
        text_width_height = _width if orient == 'vertical' else _hight
        _xaxis = x + (_width / 4.0) if orient == 'vertical' else x + scaled - (_hight * _factor)
        parent = parent.add_textbox(_xaxis, yaxis, text_width_height, text_width_height)
        data_txt_style = style.get('data', {})
        data_txt_style['color'] = data_txt_style.get('color', _color.contrast(fill))
        add_text_to_shape(parent, data_text, **data_txt_style)

    if not np.isnan(spec['target']):
        line_hight = 10000
        scaled = scale_data(spec['target'], lo, hi, factor=width)
        _width = line_hight if orient == 'horizontal' else height
        _hight = height if orient == 'horizontal' else line_hight
        yaxis = y if orient == 'horizontal' else (width - scaled) + y
        xaxis = x + scaled if orient == 'horizontal' else x
        target_line = rect(shapes, xaxis, yaxis, _width, _hight)
        fill = style.get('target', {})
        stroke = fill.get('stroke')
        fill = fill.get('fill', matplotlib.colors.to_hex(gradient(1.0)))
        rect_css(target_line, **{'fill': fill, 'stroke': stroke or fill})
        if text:
            handler = data.get('handler')
            target_text = text(spec['target'])
            parent = target_line._parent
            yaxis = yaxis - (_width / 2) if orient == 'vertical' else yaxis + (_hight / 5.0)
            text_width_height = _width if orient == 'vertical' else _hight
            parent = parent.add_textbox(xaxis, yaxis, text_width_height, text_width_height)
            target_txt_style = style.get('target', {})
            target_txt_style['color'] = target_txt_style.get('color', _color.contrast(fill))
            add_text_to_shape(parent, target_text, **target_txt_style)


def heatgrid(shape, spec, data):
    """Create a heat grid."""
    spec = copy.deepcopy(spec['heatgrid'])

    top = shape.top
    left = shape.left
    width = shape.width
    pixel_inch = 10000
    default_height = 20
    height = spec.get('cell-height', default_height) * pixel_inch
    parent = shape._parent
    shape.element.delete()

    # Loading config
    handler = data.pop('handler') if 'handler' in data else None
    for key in ['row', 'column', 'value']:
        if isinstance(spec[key], (dict,)) and 'function' in spec[key]:
            spec[key] = compile_function(spec, key, data, handler)
    # Loading data
    data = compile_function(spec, 'data', data, handler)
    data = data.sort_values(by=[spec['column']])
    rows = data[spec['row']].unique().tolist()
    columns = sorted(data[spec['column']].unique().tolist())

    left_margin = (width * spec.get('left-margin', 0.15))
    padding = spec.get('style', {}).get('padding', 10)
    padding = padding * pixel_inch

    styles = copy.deepcopy(spec.get('style', {}))

    if styles.get('gradient'):
        _min, _max = data[spec['value']].min(), data[spec['value']].max()
    # Compiling style elements if required
    for key in ['gradient', 'color', 'fill', 'font-size', 'font-family']:
        if isinstance(styles.get(key), (dict,)) and 'function' in styles[key]:
            styles[key] = compile_function(styles, key, data, handler)

    # Calculating cell's width based on config
    _width = (width - left_margin) / float(len(columns)) / pixel_inch
    _width = spec.get('cell-width', _width) * pixel_inch
    # Adding Columns to the HeatGrid.
    for idx, column in enumerate(columns):
        txt = parent.add_textbox(
            left + _width * idx + left_margin, top - height - padding, _width, height)
        add_text_to_shape(txt, '{}'.format(column), **styles)

    # Cell width
    for index, row in enumerate(rows):
        _data = data[data[spec['row']] == row].dropna()
        if len(_data) < len(columns):
            _data = pd.merge(
                _data, pd.DataFrame({spec['column']: list(columns)}),
                left_on=spec['column'], right_on=spec['column'], how='outer')
        _data = _data.sort_values(by=[spec['column']]).reset_index(drop=True)
        for _idx, _row in _data.iterrows():
            style = copy.deepcopy(styles)
            # Adding cells
            xaxis = left + (_width * _idx) + left_margin
            yaxis = top + (height * index) + padding * index
            _rect = rect(parent, xaxis, yaxis, _width - padding, height)
            # Adding color gradient to cell if gradient is True
            if style.get('gradient'):
                grad_txt = scale_data(_row[spec['value']], _min, _max)
                gradient = matplotlib.cm.get_cmap(spec['style']['gradient'])
                style['fill'] = matplotlib.colors.to_hex(gradient(grad_txt))
                style['color'] = _color.contrast(style['fill'])
            if np.isnan(_row[spec['value']]) and spec.get('na-color'):
                style['fill'] = spec.get('na-color')
                style['color'] = _color.contrast(style['fill'])

            rect_css(_rect, **style)
            # Adding text to cells if required.
            if spec.get('text'):
                spec['text'] = compile_function(spec, 'text', data, handler)
                _txt = parent.add_textbox(xaxis, yaxis, _width - padding, height)
                cell_txt = '{}'.format(_row[spec['value']])
                if callable(spec['text']):
                    cell_txt = spec['text'](_row)
                if np.isnan(_row[spec['value']]) and spec.get('na-text'):
                    cell_txt = spec.get('na-text')
                add_text_to_shape(_txt, cell_txt, **style)
        # Adding row's text in left side
        txt = parent.add_textbox(
            left, top + (height * index) + padding * index,
            _width + left_margin, height)
        add_text_to_shape(txt, row, **styles)


def css(shape, spec, data):
    """Function to modify a rectangle's property in PPT."""
    pxl_to_inch = 10000
    handler = data.pop('handler') if 'handler' in data else None
    spec = copy.deepcopy(spec['css'])
    data = compile_function(spec, 'data', data, handler)
    style = copy.deepcopy(spec.get('style', {}))
    shape_prop = {'width', 'height', 'top', 'left'}
    for prop in shape_prop:
        setprop = style.get(prop)
        if setprop:
            if not isinstance(style[prop], (dict,)):
                style[prop] = {'function': '{}'.format(style[prop]) if not isinstance(
                               style[prop], (str, six.string_types,)) else style[prop]}
            setprop = compile_function(style, prop, data, handler)
            setprop = setprop * pxl_to_inch
        else:
            setprop = getattr(shape, prop)
        setattr(shape, prop, setprop)

    _style = {}
    for key, val in style.items():
        if key not in shape_prop:
            _style[key] = val
            if isinstance(val, (dict,)):
                _style[key] = compile_function(style, key, data, handler)
            _style[key] = _style[key](data) if callable(_style[key]) else _style[key]
    rect_css(shape, **_style)


def custom_table(shape, spec, data):
    """Update an existing Table shape with data."""
    if not spec.get('custom_table', {}).get('data'):
        return
    spec = copy.deepcopy(spec['custom_table'])
    handler = data.pop('handler') if 'handler' in data else None
    data = compile_function(spec, 'data', data, handler)
    if not len(data):
        return
    data_cols = list(data.columns)
    data_cols = list(set(spec.get('columns', {}).keys() or data_cols).intersection(data_cols))
    table_properties = utils.TableProperties()
    # Extending table if required.
    table_properties.extend_table(shape, data, len(data) + 1, len(data_cols))
    # Extending table if required.
    top = shape.top
    left = shape.left
    height_margin = 5000
    two_percent = 0.02
    sixty_percent = 0.6
    twenty_percent = 0.2
    fifteen_percent = 0.15
    fourty_five_percent = 0.45
    thirty_five_percent = 0.35
    fourty_eight_percent = 0.48

    tbl_style = table_properties.get_default_css(shape)
    cell_style = table_properties.get_css(spec, data_cols, data)
    data = data.to_dict(orient='records')

    for row_num, row in enumerate(shape.table.rows):
        cols = len(row.cells._tr.tc_lst)
        # Extending cells in newly added rows.
        while cols < len(data_cols):
            row.cells._tr.add_tc()
            cols += 1
        _height = row.height + height_margin
        y = top + (_height * row_num)

        for col_num, cell in enumerate(row.cells):
            colname = data_cols[col_num]
            _width = shape.table.columns[col_num].width
            x = left + (_width * col_num)
            _d = data[row_num - 1][colname]
            if row_num > 0 and isinstance(_d, (int, float,)):
                _rect = rect(shape._parent, x + _width * twenty_percent, y + _height / 4.0,
                             _width * sixty_percent, _height / 4.0)
                circ_width = _width * fifteen_percent
                circle = shape._parent.add_shape(
                    MSO_SHAPE.OVAL, x + _width * two_percent, y + _height / 8.0,
                    circ_width, circ_width)

                masoshape = 'UP_ARROW'
                arrow_width = circ_width * fourty_five_percent
                xaxis = x + arrow_width
                yaxis = y + arrow_width
                swidth = circ_width * thirty_five_percent
                _fill = {'fill': '#119449', 'stroke': '#119449'}
                if col_num % 3:
                    _fill = {'fill': '#119449', 'stroke': '#119449'}
                elif col_num % 2:
                    masoshape = 'DOWN_ARROW'
                    _fill = {'fill': '#FEC449', 'stroke': '#FEC449'}
                else:
                    _fill = {'fill': '#D9534F', 'stroke': '#D9534F'}
                    masoshape = 'MATH_MINUS'
                    xaxis = x + circ_width * thirty_five_percent
                    swidth = circ_width * sixty_percent
                    yaxis = y + circ_width * fourty_eight_percent

                rect_css(_rect, **_fill)
                rect_css(circle, **_fill)
                arrow = shape._parent.add_shape(
                    getattr(MSO_SHAPE, masoshape), xaxis, yaxis, swidth,
                    arrow_width + arrow_width * twenty_percent)

                rect_css(arrow, **{'fill': '#ffffff', 'stroke': '#ffffff'})
            for paragraph in cell.text_frame.paragraphs:
                if not paragraph.text.strip():
                    paragraph.add_run()
                for run in paragraph.runs:
                    txt = colname if row_num == 0 else data[row_num - 1][colname]
                    run.text = '{}'.format(txt)
                    paragraph.alignment = PP_ALIGN.RIGHT
                    cellcss = {} if row_num == 0 else copy.deepcopy(cell_style.get(colname, {}))
                    txt_css = copy.deepcopy(tbl_style.get('header' if row_num == 0 else 'row', {}))
                    if row_num > 0 and 'gradient' in cellcss:
                        grad_txt = scale_data(txt, cellcss['min'], cellcss['max'])
                        gradient = matplotlib.cm.get_cmap(cellcss['gradient'])
                        cellcss['fill'] = matplotlib.colors.to_hex(gradient(grad_txt))
                    if cellcss.get('fill'):
                        txt_css['color'] = cellcss.get('color', _color.contrast(cellcss['fill']))
                        cellcss['color'] = cellcss.get('color', _color.contrast(cellcss['fill']))
                    txt_css.update(cellcss)
                    table_properties.apply_table_css(cell, paragraph, run, txt_css)


def bar_circle(shape, spec, data):
    """Function to plot bar chart with circles."""
    spec = copy.deepcopy(spec['bar_circle'])
    handler = data.pop('handler') if 'handler' in data else None
    # Getting CSS
    style = AttrDict()
    common_css = copy.deepcopy(spec.get('style', {}))
    for csskey in ['bar', 'circle']:
        _style = common_css.pop(csskey, {})
        for key, val in _style.items():
            if isinstance(val, (dict,)) and 'function' in val:
                _style[key] = compile_function(_style, key, data, handler)
        style[csskey] = _style

    for csskey in ['bar', 'circle']:
        update_css = copy.deepcopy(common_css)
        update_css.update(style[csskey])
        for key, val in update_css.items():
            if isinstance(val, (dict,)) and 'function' in val:
                update_css[key] = compile_function(update_css, key, data, handler)
        style[csskey] = update_css

    # Loading data
    bar_conf = copy.deepcopy(spec['bar'])
    bar_data = compile_function(bar_conf, 'data', data, handler)

    circle_conf = copy.deepcopy(spec['circle'])
    circle_data = compile_function(circle_conf, 'data', data, handler)

    for key in ['y', 'size', 'x', 'text']:
        for config in [bar_conf, circle_conf]:
            if isinstance(config.get(key), (dict,)) and 'function' in config[key]:
                config[key] = compile_function(config, key, data, handler)

    bary = bar_conf['y']
    circley = circle_conf['y']
    circ_size = circle_conf['size']
    bar_category = bar_conf['x']
    circle_category = circle_conf['x']
    txt_props = ['font-color', 'font-family', 'color', 'fill',
                 'font-size', 'text', 'stroke', 'opacity']
    for csskey in txt_props:
        for key in ['bar', 'circle']:
            if isinstance(style[key].get(csskey), (dict,)) and 'function' in style[key][csskey]:
                style[key][csskey] = compile_function(style[key], csskey, data, handler)

    scale = bar_data[bary].fillna(0).tolist() + circle_data[circley].fillna(0).tolist()
    ymax = max(scale)

    pixel_inch = 10000
    top = shape.top
    default_padding = 40
    default_padding_left = 4
    padding = spec.get('padding', default_padding) * pixel_inch
    left = shape.left + spec.get('padding-left', default_padding_left) * pixel_inch
    width = shape.width
    # width = width - (width * 0.20)
    ticks = width / float(len(bar_data) or 1)
    height = scale_data(bar_data[bary], 0, ymax, factor=shape.height)
    xpos = pd.np.arange(0, width, ticks)
    pos = pd.DataFrame(OrderedDict([
        ('x', xpos),
        ('y', shape.height - height),
        ('width', ticks - padding),
        ('height', height),
        ('category', bar_data[bar_category]),
        (bary, bar_data[bary])
    ]))
    # Handle negative values
    neg = pos[(pos['height'] < 0)]
    pos.loc[neg.index, 'height'] = -neg['height']
    pos.loc[neg.index, 'y'] = neg['y'] + neg['height']
    parent = shape._parent
    for index, row in pos.iterrows():
        _rect = rect(
            parent, left + row['x'], top + row['y'], row['width'], row['height'])
        txt = parent.add_textbox(
            left + row['x'] + row['width'] / 4.0,
            top + shape.height,
            row['width'] - row['width'] / 4.0, row['width'] - row['width'] / 2.0)

        _barcss = copy.deepcopy(style['bar'])
        add_text_to_shape(txt, '{}'.format(row['category']), **_barcss)

        for _barkey in txt_props:
            if callable(_barcss.get(_barkey)):
                _barcss[_barkey] = _barcss[_barkey](row)

        if bar_conf.get('text'):
            txt = parent.add_textbox(
                left + row['x'] + row['width'] / 4.0,
                top + row['y'],
                row['width'] - row['width'] / 4.0, row['width'] - row['width'] / 4.0)
            bar_txt = '{}%'.format(row[bary])
            bar_txt = bar_conf['text'](row) if callable(bar_conf['text']) else bar_txt
            add_text_to_shape(txt, bar_txt, **_barcss)

        map_circ_data = circle_data[circle_data[circle_category] == row['category']]

        for i, circle in map_circ_data.iterrows():
            y = top + scale_data(circle[circley], 0, ymax, factor=row['width'])
            y = y - (row['width'] / 2.0)
            circle_size = scale_data(
                circle[circ_size], 0, map_circ_data[circ_size].max(),
                factor=row['width'] / 2.0)
            xaxis = left + row['x'] + (row['width'] / 2.0) - (circle_size / 2.0)
            cir = parent.add_shape(
                MSO_SHAPE.OVAL, xaxis, y, circle_size, circle_size)
            circle_style = copy.deepcopy(style['circle'])
            for _key in ['fill', 'stroke', 'opacity']:
                if callable(circle_style.get(_key)):
                    circle_style[_key] = circle_style[_key](circle)
            rect_css(cir, **circle_style)
        rect_css(_rect, **_barcss)
    rect_css(shape, **{'fill': '#FBFBFB', 'stroke': '#FBFBFB'})


cmdlist = {
    'css': css,
    'text': text,
    'image': image,
    'chart': chart,
    'table': table,
    'sankey': sankey,
    'bullet': bullet,
    'replace': replace,
    'treemap': treemap,
    'heatgrid': heatgrid,
    'bar_circle': bar_circle,
    'calendarmap': calendarmap,
    'custom_table': custom_table,
}
