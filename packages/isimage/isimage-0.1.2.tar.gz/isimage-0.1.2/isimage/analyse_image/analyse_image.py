""" Command line script analyse_image
"""
# Author: Ilya Patrushev ilya.patrushev@gmail.com

# License: GPL v2.0

import os
import argparse
from isimage import Image
import matplotlib.pyplot as plt
import cv2
import numpy as np
import numpy.linalg as la
import tempfile
import shutil
import pptx


def draw_paint(img1, img2, clr):
    h = min(img1.shape[0], img2.shape[0])
    w = min(img1.shape[1], img2.shape[1])

    img1 = np.array(img1[:h, :w], dtype=float)
    img2 = np.array(img2[:h, :w], dtype=float)
    
    clr = np.array(clr, dtype=float)/256
    clr = (1 - clr)
    clr /= max(clr)
    
    img2 = img2.reshape(img2.shape[:2] + (1,))
    return (1 - ((1 - img1)*(1 - img2) + img2*clr.reshape(1, 1, 3)))

def hex2rgb(colour):
    return [int(colour[2*i + 1: 2*(i+1) + 1], 16) for i in range(3)]

def rgb2hex(colour):
    return '#%02x%02x%02x' % tuple(colour)


def main():
    parser = argparse.ArgumentParser(description=
    '''Locates an embryo in the input image and extracts the expression 
    pattern. If the output file extension is pptx the output will be a 
    PowerPoint presentation containing the original image, sketch showing 
    the outline, staining pattern and pigment pattern, and depending on 
    the verbosity level, plots showing algorithm intermediate steps.'''
    , formatter_class=argparse.ArgumentDefaultsHelpFormatter)
    parser.add_argument('-i', '--input', dest='image', metavar='image_file', type=str, nargs=1, help='path to the input image')
    parser.add_argument('-o', '--output', dest='out_name', metavar='image_file', type=str, nargs=1, help='path to save extracted pattern')
    parser.add_argument('-v', '--verbose', dest='verbose', metavar='N', type=int, default=[0], nargs=1, help='verbosity level')
    
    parser.add_argument('-m', '--max_comp', dest='max_comp', metavar='N', type=int, default=[2], nargs=1, help='maximum number of segmentation components') 
    parser.add_argument('-n', '--num_layers', dest='num_layers', metavar='N', type=int, default=[2], nargs=1, help='number of Gaussian pyramid layers to use') #  (default: [2])
    parser.add_argument('-l', '--layer_size', dest='layer_size', metavar='N', type=int, default=[200], nargs=1, help='length of the longest side of the biggest Gaussian pyramid layer to use') # (default: [200])
     
    parser.add_argument('-s', '--stain_colour', dest='stain_colour', metavar='#rrggbb', type=str, default=[rgb2hex([98, 63, 127])], nargs=1, help='stain colour') 
    parser.add_argument('-c', '--stain_confidence', dest='stain_p', metavar='R', type=float, default=[.5], nargs=1, help='confidence in staining') 
    parser.add_argument('-p', '--pigment_colour', dest='pigment_colour', metavar='#rrggbb', type=str, default=[rgb2hex([150, 140, 54])], nargs=1, help='pigment colour') 
    parser.add_argument('-C', '--pigment_confidence', dest='pigment_p', metavar='R', type=float, default=[.05], nargs=1, help='confidence in pigment presence') 
    args = parser.parse_args()
    
    if not args.image:
        print('Error: No image specified')
        args = parser.parse_args(["--help"])
        
    if not args.out_name:
        print('Error: No output file specified')
        args = parser.parse_args(["--help"])
        
    tmpdir = tempfile.mkdtemp()
    base_name = os.path.split(args.image[0])[-1]
    tmp_base = os.path.join(tmpdir, base_name)

    img = Image(args.image[0], verbose=tmp_base)
    if not img.valid:
        print args.image[0], "is not a valid image."
        return

    if args.verbose[0] < 1:
        img.verbose = False
            
    img.find_embryo(max_comp=args.max_comp[0], num_layers=args.num_layers[0], layer_size=args.layer_size[0])
        
    if not img.valid:
        print args.image[0], "is not a valid image."
        return
        
    if args.verbose[0] < 2:
        img.verbose = False

    img.extract_pattern(stain_colour=hex2rgb(args.stain_colour[0]), stain_p=args.stain_p[0], pigment_colour=hex2rgb(args.pigment_colour[0]), pigment_p=args.pigment_p[0])

    #Saving result
    if len(os.path.split(args.out_name[0])[0]) > 0 and not os.path.exists(os.path.split(args.out_name[0])[0]):
        os.makedirs(os.path.split(args.out_name[0])[0])

    saving_pptx = args.out_name[0].split('.')[-1].lower() == 'pptx'

    if not saving_pptx:
        plt.imsave(args.out_name[0], img.pattern)

    else:
        #Making sketch
        combine = np.ones_like(img.img)
        
        pattern = img.pattern.copy()
        pigment = img.pigment.copy()
        
        norm = pigment.max()
        if norm == 0:
            norm = 1
            
        exp_pigm = np.array(hex2rgb(args.pigment_colour[0]))
        exp_pigm *= 255./la.norm(hex2rgb(args.pigment_colour[0]), 2)
        
        alpha = 1 
        pigment = pigment/norm
        combine = draw_paint(combine, .75*np.float64(pigment > .5) + .33*np.float64((pigment > .25) & (pigment <= .5)), 
            np.uint8(np.round((1-alpha)*img.pigment_colour + alpha*exp_pigm))
            ) 
        
        norm = pattern.max()
        if norm == 0:
            norm = 1

        exp_stain = np.array(hex2rgb(args.stain_colour[0]))
        exp_stain *= 255./la.norm(hex2rgb(args.stain_colour[0]), 2)

        alpha = 1 
        pattern = pattern/norm
        combine = draw_paint(combine, .75*np.float64(pattern > .5) + .33*np.float64((pattern > .25) & (pattern <= .5)),  
            np.uint8(np.round((1-alpha)*img.stain_colour + alpha*exp_stain))
            ) 
        
        _, contour, _ = cv2.findContours(np.uint8(img.mask), mode=cv2.RETR_EXTERNAL, method=cv2.CHAIN_APPROX_TC89_KCOS)
        
        contour = [contour[np.argmax([cv2.contourArea(c) for c in contour])]]
        contour_simple = np.array(contour).reshape(-1, 2)
        
        cv2.drawContours(combine, [contour_simple], -1, color=(0,)*3, thickness=max(2, int(round(.002*min(combine.shape[:2])))))
            
        plt.figure()
        plt.imshow(combine)
        plt.gca().get_yaxis().set_ticks([])#(direction='out')
        plt.gca().get_xaxis().set_ticks([])#(direction='out')
        out_name = '.'.join([tmp_base, 'sketch', 'png'])
        plt.savefig(out_name)
                
        plt.figure()
        plt.imshow(img.pattern)
        plt.gca().get_yaxis().set_ticks([])#(direction='out')
        plt.gca().get_xaxis().set_ticks([])#(direction='out')
        out_name = '.'.join([tmp_base, 'expression_pattern', 'jpg'])
        plt.savefig(out_name)
        
        plt.figure()
        plt.imshow(img.pigment)
        plt.gca().get_yaxis().set_ticks([])#(direction='out')
        plt.gca().get_xaxis().set_ticks([])#(direction='out')
        out_name = '.'.join([tmp_base, 'pigment_pattern', 'jpg'])
        plt.savefig(out_name)
        
        plt.figure()
        plt.imshow(img.img)
        plt.gca().get_yaxis().set_ticks([])#(direction='out')
        plt.gca().get_xaxis().set_ticks([])#(direction='out')
        out_name = '.'.join([tmp_base, 'original_image', 'jpg'])
        plt.savefig(out_name)
            
        prs = pptx.Presentation()
        image_caption = prs.slide_layouts[1]
        
        files = {f.split('.')[-2] : f for f in os.listdir(tmpdir)}
        
        slds = ['original_image'
             , 'sketch'
             , 'expression_pattern'
             , 'pigment_pattern'
             , 'labels'
             , 'object_component'
             , 'object'
             , 'object_contour'
             , 'smoothed'
             , 'box'
             , 'stain_threshold'
             , 'pigment_threshold'
             , 'independent_axes'
             , 'proposing_colours'
             , 'classifying_colours'
            ]
            
        for s in slds:
            if s in files:
                f = files[s]
            else:
                continue
                
            slide = prs.slides.add_slide(image_caption)
            if s == 'original_image':
                slide.shapes.title.text = base_name
            else:
                slide.shapes.title.text = ' '.join(s.split('_'))
            ph = slide.shapes.placeholders[1]
            slide.shapes.add_picture('/'.join([tmpdir, f]), ph.left*1.5, ph.top*.71, height=ph.height*1.25)
        
        prs.save(args.out_name[0])
    
    shutil.rmtree(tmpdir)
    
if __name__ == '__main__':
    main()
